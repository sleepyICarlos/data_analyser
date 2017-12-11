# -*- coding: utf-8 -*-
"""
Created on Tue Dec  5 15:02:01 2017

@author: hartz
"""
from pptx import Presentation
from pptx.util import Inches


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]



def generate_presentation(title_name):
    title.text = 'Automated Powerpoint Documentation'
    subtitle.text = 'ZnSe U-I Characteristics, Resistances\n Effects from Annealing\n\nAuthor: Felix Hartz'
    prs.save(title_name)


def append_slide(figure):               # figure: filename i.e. my_pic.png
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    slide.shapes.add_picture(figure, left, top)
    
def save_pptx(presentation_name):
    prs.save(presentation_name)
    
#append_slide(fig_name)             #insert this command in the data analysis code